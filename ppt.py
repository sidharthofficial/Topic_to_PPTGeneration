import streamlit as st
from pptx import Presentation
import openai

# OpenAI API setup
openai.api_key = "YOUR_OPENAI_API_KEY"

def generate_slide_content(topic):
    """Generate slide titles and content using OpenAI."""
    prompt = (
        f"Generate an outline for a PowerPoint presentation on the topic '{topic}'.\n"
        "Include 5 slides with a title and 3-4 bullet points for each slide."
    )
    response = openai.ChatCompletion.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=200
    )
    return response['choices'][0]['message']['content']


def parse_slide_content(content):
    """Parse the OpenAI output into slide data."""
    slides = []
    lines = content.strip().split("\n")
    current_slide = None

    for line in lines:
        if line.startswith("Slide "):
            if current_slide:
                slides.append(current_slide)
            current_slide = {"title": line.split(":", 1)[1].strip(), "points": []}
        elif line.startswith("- "):
            if current_slide:
                current_slide['points'].append(line[2:].strip())

    if current_slide:
        slides.append(current_slide)

    return slides


def create_ppt(slide_data):
    """Create a PowerPoint presentation from generated data."""
    prs = Presentation()

    for slide_info in slide_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_info['title']
        content.text = "\n".join(slide_info['points'])

    output_path = "generated_presentation.pptx"
    prs.save(output_path)
    return output_path


# Streamlit App
st.title("AI-Powered PowerPoint Generator")

topic = st.text_input("Enter the topic for your presentation:")
if st.button("Generate Presentation"):
    if topic:
        with st.spinner("Generating presentation..."):
            try:
                content = generate_slide_content(topic)
                slide_data = parse_slide_content(content)
                ppt_path = create_ppt(slide_data)
                
                st.success("Presentation generated successfully!")
                with open(ppt_path, "rb") as file:
                    st.download_button(
                        label="Download Presentation",
                        data=file,
                        file_name="presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
            except Exception as e:
                st.error(f"An error occurred: {e}")
    else:
        st.warning("Please enter a topic to generate the presentation.")
